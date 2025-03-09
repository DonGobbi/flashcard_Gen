import os
import tempfile
import logging
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from youtube_transcript_api import YouTubeTranscriptApi
from groq import Groq
from services.ai_service import AIService
from services.auth_service import AuthService
from functools import wraps
from fpdf import FPDF
import genanki
import json
import dotenv
from docx import Document
from pptx import Presentation
import csv
import io
import PyPDF2

# Load environment variables
dotenv.load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))

app = Flask(__name__)
CORS(app, resources={
    r"/api/*": {
        "origins": ["http://localhost:3000"],
        "methods": ["GET", "POST", "PUT", "DELETE", "OPTIONS"],
        "allow_headers": ["Content-Type", "Authorization", "Accept"],
        "expose_headers": ["Content-Type"],
        "max_age": 3600,
        "supports_credentials": True
    }
})

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize services
ai_service = AIService()
auth_service = AuthService()
groq_client = Groq(api_key=os.getenv('GROQ_API_KEY'))

def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = None
        if 'Authorization' in request.headers:
            token = request.headers['Authorization'].replace('Bearer ', '')
        
        if not token:
            return jsonify({'error': 'Token is missing'}), 401

        try:
            email = auth_service.verify_token(token)
            if not email:
                return jsonify({'error': 'Invalid token'}), 401
        except Exception as e:
            return jsonify({'error': 'Invalid token'}), 401

        return f(*args, **kwargs)
    return decorated

def extract_text_from_file(file):
    filename = file.filename.lower()
    content = file.read()
    
    try:
        if filename.endswith('.txt') or filename.endswith('.md'):
            return content.decode('utf-8')
            
        elif filename.endswith('.pdf'):
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = ''
            for page in pdf_reader.pages:
                text += page.extract_text() + '\n'
            return text
            
        elif filename.endswith('.docx'):
            doc = Document(io.BytesIO(content))
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            
        elif filename.endswith('.pptx'):
            prs = Presentation(io.BytesIO(content))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
            
        elif filename.endswith('.csv'):
            csv_content = content.decode('utf-8').splitlines()
            reader = csv.reader(csv_content)
            return '\n'.join([' '.join(row) for row in reader])
            
        else:
            raise ValueError('Unsupported file format')
            
    except Exception as e:
        raise ValueError(f'Error processing file: {str(e)}')

def extract_video_id(url):
    if 'youtu.be' in url:
        return url.split('/')[-1]
    if 'v=' in url:
        return url.split('v=')[1].split('&')[0]
    return None

def generate_flashcards_from_text(text, num_cards=5):
    prompt = f"""Given the following text, generate {num_cards} flashcards in a question-answer format. 
    Make the questions clear and concise, and ensure the answers are accurate based on the content.
    Format the output as a list of dictionaries with 'question' and 'answer' keys.
    Text: {text}"""

    chat_completion = groq_client.chat.completions.create(
        messages=[
            {
                "role": "system",
                "content": "You are a helpful assistant that creates educational flashcards."
            },
            {
                "role": "user",
                "content": prompt
            }
        ],
        model="mixtral-8x7b-32768",
        temperature=0.7,
        max_tokens=2048,
    )

    try:
        response_text = chat_completion.choices[0].message.content
        import ast
        flashcards = ast.literal_eval(response_text)
        return flashcards
    except Exception as e:
        logger.error(f"Error parsing flashcards: {e}")
        return []

@app.route('/api/auth/register', methods=['POST'])
def register():
    try:
        data = request.get_json()
        if not data or 'email' not in data or 'password' not in data:
            return jsonify({'error': 'Email and password are required'}), 400

        result = auth_service.register_user(data['email'], data['password'])
        return jsonify(result)
    except ValueError as e:
        return jsonify({'error': str(e)}), 400
    except Exception as e:
        logger.error(f"Error in register: {str(e)}")
        return jsonify({'error': 'Registration failed'}), 500

@app.route('/api/auth/login', methods=['POST'])
def login():
    try:
        data = request.get_json()
        if not data or 'email' not in data or 'password' not in data:
            return jsonify({'error': 'Email and password are required'}), 400

        result = auth_service.login_user(data['email'], data['password'])
        return jsonify(result)
    except ValueError as e:
        return jsonify({'error': str(e)}), 400
    except Exception as e:
        logger.error(f"Error in login: {str(e)}")
        return jsonify({'error': 'Login failed'}), 500

@app.route('/api/youtube', methods=['POST'])
@token_required
def process_youtube():
    try:
        data = request.get_json()
        if not data:
            return jsonify({'error': 'Invalid JSON data'}), 400
            
        url = data.get('url')
        num_cards = int(data.get('num_cards', 5))
        
        if not url:
            return jsonify({'error': 'No URL provided'}), 400

        video_id = extract_video_id(url)
        if not video_id:
            return jsonify({'error': 'Invalid YouTube URL'}), 400

        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
            text = ' '.join([entry['text'] for entry in transcript])
            flashcards = generate_flashcards_from_text(text, num_cards)
            
            return jsonify({'flashcards': flashcards})
        except Exception as e:
            return jsonify({'error': f'Failed to get transcript: {str(e)}'}), 400

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/upload', methods=['POST'])
@token_required
def process_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400

        file = request.files['file']
        num_cards = int(request.form.get('num_cards', 5))
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400

        try:
            text = extract_text_from_file(file)
            flashcards = generate_flashcards_from_text(text, num_cards)
            return jsonify({'flashcards': flashcards})
        except ValueError as e:
            return jsonify({'error': str(e)}), 400

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/improve', methods=['POST'])
@token_required
def improve_flashcard():
    try:
        data = request.get_json()
        if not data or 'flashcard' not in data:
            return jsonify({'error': 'No flashcard provided'}), 400
            
        improved_flashcard = ai_service.improve_flashcard(data['flashcard'])
        return jsonify({'flashcard': improved_flashcard})
    except Exception as e:
        logger.error(f"Error in improve_flashcard: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/translate', methods=['POST'])
@token_required
def translate_flashcard():
    try:
        data = request.get_json()
        if not data or 'flashcard' not in data or 'target_language' not in data:
            return jsonify({'error': 'Missing flashcard or target language'}), 400
            
        translated_flashcard = ai_service.translate_flashcard(
            data['flashcard'],
            data['target_language']
        )
        return jsonify({'flashcard': translated_flashcard})
    except Exception as e:
        logger.error(f"Error in translate_flashcard: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/export/pdf', methods=['POST'])
@token_required
def export_pdf():
    try:
        data = request.get_json()
        if not data or 'flashcards' not in data:
            return jsonify({'error': 'No flashcards provided'}), 400

        customization = data.get('customization', {})
        style = customization.get('style', 'Classic')
        
        # Create PDF with customization
        pdf = FPDF()
        pdf.add_page()
        
        # Set font based on customization
        font_family = customization.get('font', 'Arial')
        font_size = int(customization.get('fontSize', 12))
        
        # Add title with custom styling
        pdf.set_font(font_family if font_family in ['Arial', 'Times'] else 'Arial', 
                    style='B', 
                    size=font_size + 6)
        pdf.set_text_color(33, 150, 243)  # Primary blue color
        pdf.cell(0, 20, txt="My Flashcards", ln=1, align='C')
        pdf.ln(10)

        # Style-specific settings
        if style == 'Modern':
            bg_color = (245, 245, 245)
            border_color = (33, 150, 243)
            question_color = (33, 150, 243)
            answer_color = (85, 85, 85)
        elif style == 'Minimalist':
            bg_color = (255, 255, 255)
            border_color = (200, 200, 200)
            question_color = (0, 0, 0)
            answer_color = (85, 85, 85)
        elif style == 'Colorful':
            bg_color = (240, 248, 255)
            border_color = (33, 150, 243)
            question_color = (156, 39, 176)
            answer_color = (0, 150, 136)
        else:  # Classic
            bg_color = (250, 250, 250)
            border_color = (180, 180, 180)
            question_color = (33, 33, 33)
            answer_color = (85, 85, 85)

        # Add flashcards with custom styling
        for i, card in enumerate(data['flashcards'], 1):
            # Card container
            pdf.set_fill_color(*bg_color)
            pdf.set_draw_color(*border_color)
            pdf.rect(10, pdf.get_y(), 190, 0, 'S')
            pdf.ln(5)

            # Card number
            pdf.set_font(font_family if font_family in ['Arial', 'Times'] else 'Arial', 
                        style='B', 
                        size=font_size)
            pdf.set_text_color(*question_color)
            pdf.cell(0, 10, txt=f"Card {i}", ln=1, align='L')
            
            # Question
            pdf.set_font(font_family if font_family in ['Arial', 'Times'] else 'Arial', 
                        style='B', 
                        size=font_size)
            pdf.set_text_color(*question_color)
            pdf.multi_cell(0, 10, txt="Question:", fill=True)
            pdf.set_font(font_family if font_family in ['Arial', 'Times'] else 'Arial', 
                        size=font_size)
            pdf.multi_cell(0, 10, txt=card['question'])
            pdf.ln(5)
            
            # Answer
            pdf.set_font(font_family if font_family in ['Arial', 'Times'] else 'Arial', 
                        style='B', 
                        size=font_size)
            pdf.set_text_color(*answer_color)
            pdf.multi_cell(0, 10, txt="Answer:", fill=True)
            pdf.set_font(font_family if font_family in ['Arial', 'Times'] else 'Arial', 
                        size=font_size)
            pdf.multi_cell(0, 10, txt=card['answer'])
            pdf.ln(10)

        # Save PDF to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
            pdf.output(tmp.name)
            return send_file(
                tmp.name,
                mimetype='application/pdf',
                as_attachment=True,
                download_name='flashcards.pdf'
            )

    except Exception as e:
        logger.error(f"Error in export_pdf: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/export/anki', methods=['POST'])
@token_required
def export_anki():
    try:
        data = request.get_json()
        if not data or 'flashcards' not in data:
            return jsonify({'error': 'No flashcards provided'}), 400

        # Create a new deck
        deck_name = data.get('deckName', 'My Flashcards')
        deck_id = hash(deck_name + str(os.urandom(32)))  # Generate a random deck ID
        deck = genanki.Deck(deck_id, deck_name)

        # Define the note model (template for cards)
        model = genanki.Model(
            1607392319,  # Random model ID
            'Simple Model',
            fields=[
                {'name': 'Question'},
                {'name': 'Answer'},
            ],
            templates=[
                {
                    'name': 'Card 1',
                    'qfmt': '{{Question}}',
                    'afmt': '{{FrontSide}}<hr id="answer">{{Answer}}',
                },
            ],
            css="""
            .card {
                font-family: arial;
                font-size: 20px;
                text-align: center;
                color: black;
                background-color: white;
            }
            """
        )

        # Add notes (cards) to the deck
        for card in data['flashcards']:
            note = genanki.Note(
                model=model,
                fields=[card['question'], card['answer']]
            )
            deck.add_note(note)

        # Create a package containing the deck
        package = genanki.Package(deck)
        
        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.apkg') as tmp:
            package.write_to_file(tmp.name)
            return send_file(
                tmp.name,
                mimetype='application/apkg',
                as_attachment=True,
                download_name='flashcards.apkg'
            )

    except Exception as e:
        logger.error(f"Error in export_anki: {str(e)}")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)