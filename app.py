from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import os
from werkzeug.utils import secure_filename
from werkzeug.exceptions import HTTPException
import docx
from pptx import Presentation
import csv
import pandas as pd
from services.ai_service import AIService
import asyncio
from functools import wraps
import uuid
from dotenv import load_dotenv
from config import Config
import shutil

# Load environment variables from .env file
load_dotenv()

app = Flask(__name__, static_folder='frontend/build')
CORS(app)

# Initialize configuration
Config.init_app(app)

# Configure upload folder
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'temp')
# Remove temp directory if it exists and recreate it
if os.path.exists(UPLOAD_FOLDER):
    shutil.rmtree(UPLOAD_FOLDER)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Configure allowed file extensions
ALLOWED_EXTENSIONS = {'txt', 'docx', 'pptx', 'csv'}

# Configure max content length (50MB)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024

# Initialize AI service
ai_service = AIService()

def async_route(f):
    @wraps(f)
    async def wrapped(*args, **kwargs):
        return await f(*args, **kwargs)
    return wrapped

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def process_docx(file_path):
    """Extract content from DOCX files, preserving structure."""
    try:
        if not os.path.exists(file_path):
            raise ValueError(f"File not found: {file_path}")
            
        if not os.access(file_path, os.R_OK):
            raise ValueError(f"File not readable: {file_path}")
            
        doc = docx.Document(file_path)
        content = []
        
        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)
                
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        content.append(cell.text)
        
        if not content:
            raise ValueError("No content found in the document")
            
        return '\n\n'.join(content)
        
    except docx.opc.exceptions.PackageNotFoundError:
        raise ValueError("Invalid or corrupted DOCX file")
    except Exception as e:
        app.logger.error(f"Error processing DOCX file: {str(e)}")
        raise ValueError(f"Failed to process DOCX file: {str(e)}")

def process_pptx(file_path):
    """Extract content from PPTX files."""
    try:
        prs = Presentation(file_path)
        content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
        return '\n\n'.join(content)
    except Exception as e:
        app.logger.error(f"Error processing PPTX file: {str(e)}")
        raise

def process_txt(file_path):
    """Process text files."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        app.logger.error(f"Error processing TXT file: {str(e)}")
        raise

def process_csv(file_path):
    """Process CSV files."""
    try:
        df = pd.read_csv(file_path)
        return '\n'.join([
            ', '.join(map(str, row)) 
            for row in df.values
        ])
    except Exception as e:
        app.logger.error(f"Error processing CSV file: {str(e)}")
        raise

@app.errorhandler(Exception)
def handle_error(error):
    """Global error handler to ensure JSON responses"""
    error_message = str(error)
    error_type = error.__class__.__name__
    
    app.logger.error(f"Error Type: {error_type}")
    app.logger.error(f"Error Message: {error_message}")
    
    if isinstance(error, HTTPException):
        status_code = error.code
        message = error.description
    else:
        status_code = 500
        message = f"{error_type}: {error_message}"
    
    response = {
        "error": message,
        "type": error_type
    }
    return jsonify(response), status_code

# Serve React App
@app.route('/', defaults={'path': ''})
@app.route('/<path:path>')
def serve(path):
    if path != "" and os.path.exists(app.static_folder + '/' + path):
        return send_from_directory(app.static_folder, path)
    else:
        return send_from_directory(app.static_folder, 'index.html')

@app.route('/api/upload', methods=['POST'])
@async_route
async def upload_file():
    temp_path = None
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
            
        if not allowed_file(file.filename):
            return jsonify({"error": "File type not supported. Please upload a .txt, .docx, .pptx, or .csv file"}), 400

        num_cards = request.form.get('num_cards', type=int, default=5)
        subject = request.form.get('subject', '')

        # Create a unique filename to avoid conflicts
        filename = secure_filename(str(uuid.uuid4()) + '_' + file.filename)
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # Ensure temp directory exists
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        
        # Save file temporarily
        try:
            file.save(temp_path)
            if not os.path.exists(temp_path):
                raise ValueError("Failed to save file")
            if not os.access(temp_path, os.R_OK):
                raise ValueError("File saved but not readable")
            file_size = os.path.getsize(temp_path)
            if file_size == 0:
                raise ValueError("File saved but is empty")
            app.logger.info(f"File saved successfully: {temp_path} ({file_size} bytes)")
        except Exception as e:
            raise ValueError(f"Failed to save file: {str(e)}")

        try:
            # Extract text from file
            if file.filename.endswith('.docx'):
                text = process_docx(temp_path)
            elif file.filename.endswith('.pptx'):
                text = process_pptx(temp_path)
            elif file.filename.endswith('.txt'):
                text = process_txt(temp_path)
            elif file.filename.endswith('.csv'):
                text = process_csv(temp_path)
            else:
                raise ValueError("Unsupported file format")

            if not text or not text.strip():
                raise ValueError("No text content found in file")

            app.logger.info(f"Successfully extracted text from {filename}")

            # Generate flashcards using AI
            flashcards = await ai_service.generate_flashcards(text, num_cards, subject)
            if not flashcards:
                raise ValueError("Failed to generate flashcards")
            
            app.logger.info(f"Successfully generated {len(flashcards)} flashcards")
            return jsonify({"flashcards": flashcards})

        except Exception as e:
            app.logger.error(f"Error processing file {filename}: {str(e)}")
            raise

    except Exception as e:
        error_msg = str(e)
        app.logger.error(f"Error in upload_file: {error_msg}")
        return jsonify({"error": error_msg}), 500

    finally:
        # Clean up temp file
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
                app.logger.info(f"Cleaned up temp file: {temp_path}")
            except Exception as e:
                app.logger.error(f"Error cleaning up temp file {temp_path}: {str(e)}")

@app.route('/api/improve', methods=['POST'])
@async_route
async def improve_flashcard():
    try:
        data = request.get_json()
        if not data or 'question' not in data or 'answer' not in data:
            return jsonify({"error": "Missing question or answer"}), 400

        improved = await ai_service.improve_flashcard(data['question'], data['answer'])
        return jsonify(improved)

    except Exception as e:
        app.logger.error(f"Error in improve_flashcard: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/translate', methods=['POST'])
@async_route
async def translate_flashcard():
    try:
        data = request.get_json()
        if not data or 'question' not in data or 'answer' not in data or 'target_language' not in data:
            return jsonify({"error": "Missing required fields"}), 400
        
        translated = await ai_service.translate_flashcard(
            data['question'], 
            data['answer'],
            data['target_language']
        )
        return jsonify(translated)

    except Exception as e:
        app.logger.error(f"Error translating flashcard: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({"status": "healthy"}), 200

if __name__ == '__main__':
    import asyncio
    from hypercorn.config import Config
    from hypercorn.asyncio import serve

    config = Config()
    config.bind = ["localhost:5000"]
    config.use_reloader = True
    
    asyncio.run(serve(app, config))
